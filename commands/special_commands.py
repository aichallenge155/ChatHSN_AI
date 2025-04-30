import discord
from discord.ext import commands
import os
import requests
from dotenv import load_dotenv
from database import add_user, update_last_interaction, log_interaction
from ai import find_citation, define_term, generate_studyplan
from googletrans import Translator
from PyPDF2 import PdfReader, PdfWriter
from docx import Document
import openpyxl
import pptx
from io import BytesIO
import tempfile
from pdf2docx import Converter
from docx2pdf import convert
import re
import asyncio

load_dotenv()

class SpecialCommands(commands.Cog):
    def __init__(self, bot):
        self.bot = bot
        self.stability_api_key = os.getenv("STABILITY_API_KEY")
        self.translator = Translator()
        # Define supported conversions for interactive menu
        self.supported_conversions = [
            {"from": "pdf", "to": "docx", "description": "PDF → DOCX (Word document)"},
            {"from": "docx", "to": "pdf", "description": "DOCX → PDF"},
            {"from": "pdf", "to": "txt", "description": "PDF → TXT (Text file)"},
            {"from": "xlsx", "to": "csv", "description": "XLSX → CSV (Comma separated values)"}
        ]

    @commands.command(name="cite")
    async def cite(self, ctx, *, topic=None):
        """Find a citation for a topic."""
        if not topic:
            await ctx.send("Please provide a topic to find a citation for.")
            return
        try:
            # Log user interaction in the database
            add_user(str(ctx.author.id), ctx.author.name)
            update_last_interaction(str(ctx.author.id))
            log_interaction(str(ctx.author.id), f"cite: {topic}")
        except Exception as e:
            await ctx.send(f"Error logging interaction: {e}")
            print(f"Error logging interaction: {e}")
            return

        async with ctx.typing():
            try:
                citation = await find_citation(topic)
                if citation is None:
                    await ctx.send("No citation found for this topic.")
                    return
                    
                embed = discord.Embed(
                    title=f"Citation for: {topic}",
                    description=citation,
                    color=discord.Color.gold()
                )
                await ctx.send(embed=embed)
            except Exception as e:
                await ctx.send(f"Error retrieving citation: {e}")
                print(f"Error retrieving citation: {e}")

    @commands.command(name="define")
    async def define(self, ctx, *, term=None):
        """Define a term."""
        if not term:
            await ctx.send("Please provide a term to define.")
            return
        try:
            # Log user interaction in the database
            add_user(str(ctx.author.id), ctx.author.name)
            update_last_interaction(str(ctx.author.id))
            log_interaction(str(ctx.author.id), f"define: {term}")
        except Exception as e:
            await ctx.send(f"Error logging interaction: {e}")
            print(f"Error logging interaction: {e}")
            return

        async with ctx.typing():
            try:
                definition = await define_term(term)
                if definition is None:
                    await ctx.send("No definition found for this term.")
                    return
                    
                embed = discord.Embed(
                    title=f"Definition: {term}",
                    description=definition,
                    color=discord.Color.blue()
                )
                await ctx.send(embed=embed)
            except Exception as e:
                await ctx.send(f"Error retrieving definition: {e}")
                print(f"Error retrieving definition: {e}")

    @commands.command(name="studyplan")
    async def studyplan(self, ctx, *, topic=None):
        """Generate a motivational study plan for a topic."""
        if not topic:
            await ctx.send("📚 Zəhmət olmasa plan üçün mövzu və ya fənn daxil et. Məsələn: `!studyplan ingilis dili imtahanı`")
            return
        
        # Log user interaction in database
        try:
            add_user(str(ctx.author.id), ctx.author.name)
            update_last_interaction(str(ctx.author.id))
            log_interaction(str(ctx.author.id), f"studyplan: {topic}")
        except Exception as e:
            await ctx.send(f"Error logging interaction: {e}")
            print(f"Error logging interaction: {e}")
            return
        
        # Generate study plan
        async with ctx.typing():
            try:
                plan = await generate_studyplan(topic)
                if plan is None:
                    await ctx.send("Study plan generation failed. Please try again later.")
                    return
                
                # Split plan into chunks if it exceeds Discord's embed limit
                MAX_EMBED_DESCRIPTION_LENGTH = 4000
                
                if len(plan) <= MAX_EMBED_DESCRIPTION_LENGTH:
                    # If plan is short, send as a single embed
                    embed = discord.Embed(
                        title=f"📘 Study Plan: {topic.title()}",
                        description=plan,
                        color=discord.Color.green()
                    )
                    await ctx.send(embed=embed)
                else:
                    # If plan is long, split into multiple parts
                    await ctx.send(f"📘 **Study Plan for {topic.title()}**")
                    
                    # Split the plan into parts
                    parts = []
                    for i in range(0, len(plan), MAX_EMBED_DESCRIPTION_LENGTH):
                        parts.append(plan[i:i + MAX_EMBED_DESCRIPTION_LENGTH])
                    
                    # Send each part as a separate embed
                    for i, part in enumerate(parts):
                        part_embed = discord.Embed(
                            title=f"Part {i+1}/{len(parts)}",
                            description=part,
                            color=discord.Color.green()
                        )
                        await ctx.send(embed=part_embed)
                    
            except Exception as e:
                await ctx.send(f"Tədris planı yaradılarkən xəta baş verdi: {str(e)}")
                print(f"Error generating study plan: {e}")

    @commands.command(name="image")
    async def generate_image(self, ctx, *, prompt: str):
        """Generate an image based on a prompt."""
        if not self.stability_api_key:
            await ctx.send("❌ API açarı tapılmadı.")
            return

        # Log user interaction
        try:
            add_user(str(ctx.author.id), ctx.author.name)
            update_last_interaction(str(ctx.author.id))
            log_interaction(str(ctx.author.id), f"image: {prompt}")
        except Exception as e:
            await ctx.send(f"Error logging interaction: {e}")
            print(f"Error logging interaction: {e}")
            return

        # Show typing animation while creating the image
        async with ctx.typing():
            await ctx.send(f"🎨 Şəkil yaradılır: **{prompt}**...")

            # Translate prompt to English if needed
            try:
                translated_prompt = self.translator.translate(prompt, src='auto', dest='en').text
                print(f"Original Prompt: {prompt}")
                print(f"Translated Prompt: {translated_prompt}")
            except Exception as e:
                await ctx.send(f"Translation error: {e}")
                print(f"Translation error: {e}")
                # Continue with original prompt if translation fails
                translated_prompt = prompt

            # Create a temporary file for the image
            temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            temp_filename = temp_file.name
            temp_file.close()
            
            try:
                # Send POST request to Stability AI API
                response = requests.post(
                    "https://api.stability.ai/v1/generation/stable-diffusion-xl-1024-v1-0/text-to-image",
                    headers={
                        "Authorization": f"Bearer {self.stability_api_key}",
                        "Accept": "image/png",
                        "Content-Type": "application/json"
                    },
                    json={
                        "text_prompts": [{"text": translated_prompt}],
                        "cfg_scale": 7,
                        "height": 1024,
                        "width": 1024,
                        "samples": 1,
                        "steps": 30,
                    }
                )

                if response.status_code == 200:
                    # Save the image to the temporary file
                    with open(temp_filename, "wb") as f:
                        f.write(response.content)
                    
                    # Send the image to Discord
                    await ctx.send(file=discord.File(temp_filename))
                else:
                    await ctx.send(f"API error (Status {response.status_code}): {response.text}")
                    print(f"API error: {response.text}")
            except Exception as e:
                await ctx.send(f"Error generating image: {e}")
                print(f"Error generating image: {e}")
            finally:
                # Clean up the temporary file
                try:
                    os.unlink(temp_filename)
                except Exception as e:
                    print(f"Error deleting temporary file: {e}")

    @commands.command(name="analyse")
    async def read_file(self, ctx):
        """Read contents of files (PDF, DOCX, PPTX, XLSX)."""
        if len(ctx.message.attachments) == 0:
            await ctx.send("❌ Zəhmət olmasa, oxumaq istədiyiniz faylı əlavə edin.")
            return
        
        # Log user interaction
        try:
            add_user(str(ctx.author.id), ctx.author.name)
            update_last_interaction(str(ctx.author.id))
            log_interaction(str(ctx.author.id), "readfile: file")
        except Exception as e:
            await ctx.send(f"Error logging interaction: {e}")
            print(f"Error logging interaction: {e}")
            return
        
        file = ctx.message.attachments[0]
        file_bytes = await file.read()
        file_type = file.content_type
        
        try:
            # Process file types based on content type
            if "pdf" in file_type:
                await self.read_pdf(ctx, file_bytes, file.filename)
            elif "word" in file_type:
                await self.read_docx(ctx, file_bytes, file.filename)
            elif "presentation" in file_type:
                await self.read_pptx(ctx, file_bytes, file.filename)
            elif "spreadsheet" in file_type:
                await self.read_xlsx(ctx, file_bytes, file.filename)
            else:
                await ctx.send(f"❌ Dəstəklənməyən fayl növü: {file_type}")
        except Exception as e:
            await ctx.send(f"❌ Fayl oxunarkən xəta baş verdi: {e}")
            print(f"Error reading file: {e}")

    async def read_pdf(self, ctx, file_bytes, filename):
        """Read a PDF file."""
        try:
            reader = PdfReader(BytesIO(file_bytes))
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            
            # Split the text into manageable chunks
            await self.send_text_in_chunks(ctx, text, f"📄 PDF faylının məzmunu ({filename}):")
        except Exception as e:
            await ctx.send(f"❌ PDF oxunarkən xəta baş verdi: {e}")
            print(f"Error reading PDF: {e}")

    async def read_docx(self, ctx, file_bytes, filename):
        """Read a DOCX file."""
        try:
            doc = Document(BytesIO(file_bytes))
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Split the text into manageable chunks
            await self.send_text_in_chunks(ctx, text, f"📄 DOCX faylının məzmunu ({filename}):")
        except Exception as e:
            await ctx.send(f"❌ DOCX oxunarkən xəta baş verdi: {e}")
            print(f"Error reading DOCX: {e}")
    
    async def read_pptx(self, ctx, file_bytes, filename):
        """Read a PPTX file."""
        try:
            presentation = pptx.Presentation(BytesIO(file_bytes))
            text = ""
            for i, slide in enumerate(presentation.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            # Split the text into manageable chunks
            await self.send_text_in_chunks(ctx, text, f"📄 PPTX faylının məzmunu ({filename}):")
        except Exception as e:
            await ctx.send(f"❌ PPTX oxunarkən xəta baş verdi: {e}")
            print(f"Error reading PPTX: {e}")

    async def read_xlsx(self, ctx, file_bytes, filename):
        """Read an XLSX file."""
        try:
            workbook = openpyxl.load_workbook(BytesIO(file_bytes))
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"--- Sheet: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    text += " | ".join(str(cell) if cell is not None else "" for cell in row) + "\n"
                text += "\n"
            
            # Split the text into manageable chunks
            await self.send_text_in_chunks(ctx, text, f"📄 XLSX faylının məzmunu ({filename}):")
        except Exception as e:
            await ctx.send(f"❌ XLSX oxunarkən xəta baş verdi: {e}")
            print(f"Error reading XLSX: {e}")

    async def send_text_in_chunks(self, ctx, text, header):
        """Send text in chunks to avoid Discord's message length limit."""
        MAX_MESSAGE_LENGTH = 1900  # Discord message limit is 2000 chars

        # Send the header first
        await ctx.send(header)
        
        # Split the text into chunks and send each chunk
        if len(text) <= MAX_MESSAGE_LENGTH:
            await ctx.send(f"```\n{text}\n```")
        else:
            chunks = []
            for i in range(0, len(text), MAX_MESSAGE_LENGTH):
                chunks.append(text[i:i + MAX_MESSAGE_LENGTH])
            
            for i, chunk in enumerate(chunks):
                await ctx.send(f"```\nPart {i+1}/{len(chunks)}:\n{chunk}\n```")

    @commands.command(name="convert")
    async def convert(self, ctx, from_format=None, to_format=None):
        """
        Convert a file from one format to another format
        Interactive menu if no formats specified
        """
        # Check if user attached a file
        if len(ctx.message.attachments) == 0:
            await ctx.send("❌ Zəhmət olmasa, çevirmək istədiyiniz faylı əlavə edin.")
            return
        
        # Log user interaction
        try:
            add_user(str(ctx.author.id), ctx.author.name)
            update_last_interaction(str(ctx.author.id))
            log_interaction(str(ctx.author.id), "convert: started")
        except Exception as e:
            await ctx.send(f"Error logging interaction: {e}")
            print(f"Error logging interaction: {e}")
            return
            
        attached_file = ctx.message.attachments[0]
        file_extension = attached_file.filename.split('.')[-1].lower()
        
        # If formats not specified, show interactive menu
        if from_format is None or to_format is None:
            # Filter conversion options based on the attached file type
            available_conversions = [
                conversion for conversion in self.supported_conversions 
                if conversion["from"] == file_extension
            ]
            
            if not available_conversions:
                await ctx.send(f"❌ Bu fayl növü ({file_extension}) üçün çevirmə əməliyyatı mövcud deyil.")
                return
                
            # Create the interactive menu embed
            embed = discord.Embed(
                title="🔄 Fayl çevirməsi", 
                description=f"Aşağıdakı seçimlər üçün reaksiya bildirin:\n({attached_file.filename})",
                color=discord.Color.blue()
            )
            
            # Add options with emoji numbers
            emoji_numbers = ["1️⃣", "2️⃣", "3️⃣", "4️⃣", "5️⃣", "6️⃣", "7️⃣", "8️⃣", "9️⃣"]
            conversion_options = {}
            
            # Make sure we don't exceed available options
            available_options = min(len(emoji_numbers), len(available_conversions))
            
            for i in range(available_options):
                emoji = emoji_numbers[i]
                conversion = available_conversions[i]
                embed.add_field(
                    name=f"{emoji} {conversion['description']}", 
                    value=f"`{conversion['from']}` ➡️ `{conversion['to']}`", 
                    inline=False
                )
                conversion_options[emoji] = conversion
                
            embed.set_footer(text="30 saniyə ərzində bir seçim edin")
            menu_message = await ctx.send(embed=embed)
            
            # Add reaction options
            for emoji in conversion_options.keys():
                await menu_message.add_reaction(emoji)
                
            # Wait for user reaction
            def check(reaction, user):
                return (
                    user == ctx.author and 
                    str(reaction.emoji) in conversion_options.keys() and 
                    reaction.message.id == menu_message.id
                )
                
            try:
                reaction, user = await self.bot.wait_for('reaction_add', timeout=30.0, check=check)
                selected_conversion = conversion_options[str(reaction.emoji)]
                from_format = selected_conversion["from"]
                to_format = selected_conversion["to"]
                
                # Update the menu to show selection
                embed.title = f"🔄 Seçilmiş çevirmə: {selected_conversion['description']}"
                embed.color = discord.Color.green()
                await menu_message.edit(embed=embed)
                
            except asyncio.TimeoutError:
                embed.description = "❌ Vaxt bitdi. Zəhmət olmasa yenidən cəhd edin."
                embed.color = discord.Color.red()
                await menu_message.edit(embed=embed)
                return
        
        # Verify that the file extension matches the from_format
        if file_extension != from_format:
            await ctx.send(f"❌ Faylın növü ({file_extension}) seçilmiş mənbə formatı ilə ({from_format}) uyğun gəlmir.")
            return

        # Let user know we're working on it
        async with ctx.typing():
            await ctx.send(f"⏳ {attached_file.filename} faylı {to_format} formatına çevrilir...")
            
            # Create a temporary directory for processing
            temp_dir = tempfile.mkdtemp()
            input_path = os.path.join(temp_dir, attached_file.filename)
            output_filename = f"{os.path.splitext(attached_file.filename)[0]}.{to_format}"
            output_path = os.path.join(temp_dir, output_filename)
            
            try:
                # Download the file
                file_content = await attached_file.read()
                with open(input_path, 'wb') as f:
                    f.write(file_content)
                
                # Do the conversion based on the formats
                if from_format == "pdf" and to_format == "docx":
                    await self.convert_pdf_to_docx(input_path, output_path)
                elif from_format == "docx" and to_format == "pdf":
                    await self.convert_docx_to_pdf(input_path, output_path)
                elif from_format == "pdf" and to_format == "txt":
                    await self.convert_pdf_to_txt(input_path, output_path)
                elif from_format == "xlsx" and to_format == "csv":
                    await self.convert_xlsx_to_csv(input_path, output_path)
                else:
                    await ctx.send(f"❌ Dəstəklənməyən çevirmə: {from_format} ➡️ {to_format}")
                    return

                # Send the converted file
                await ctx.send("✅ Çevirmə tamamlandı!", file=discord.File(output_path))
                
                # Log successful conversion
                log_interaction(str(ctx.author.id), f"convert: {from_format} to {to_format} (successful)")
                
            except Exception as e:
                await ctx.send(f"❌ Çevirmə zamanı xəta baş verdi: {str(e)}")
                print(f"Conversion error: {e}")
                log_interaction(str(ctx.author.id), f"convert: {from_format} to {to_format} (failed: {str(e)})")
            finally:
                # Cleanup temporary files
                try:
                    if os.path.exists(input_path):
                        os.remove(input_path)
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    os.rmdir(temp_dir)
                except Exception as e:
                    print(f"Error cleaning up temp files: {e}")

    async def convert_pdf_to_docx(self, pdf_path, docx_path):
        """Convert PDF file to DOCX format."""
        # Using pdf2docx library
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()

    async def convert_docx_to_pdf(self, docx_path, pdf_path):
        """Convert DOCX file to PDF format."""
        # Using docx2pdf library
        convert(docx_path, pdf_path)

    async def convert_pdf_to_txt(self, pdf_path, txt_path):
        """Convert PDF file to plain text."""
        # Using PyPDF2
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            with open(txt_path, 'w', encoding='utf-8') as output_file:
                for page_num in range(len(reader.pages)):
                    text = reader.pages[page_num].extract_text()
                    output_file.write(text)

    async def convert_xlsx_to_csv(self, xlsx_path, csv_path):
        """Convert XLSX file to CSV format."""
        wb = openpyxl.load_workbook(xlsx_path)
        sheet = wb.active
        
        with open(csv_path, 'w', newline='', encoding='utf-8') as f:
            for row in sheet.rows:
                values = [str(cell.value) if cell.value is not None else '' for cell in row]
                f.write(','.join(values) + '\n')

async def setup(bot):
    await bot.add_cog(SpecialCommands(bot))